"""
Domain Services

비즈니스 로직을 담당하는 도메인 서비스입니다.
단일책임원칙과 함수형 프로그래밍을 적용합니다.
모든 메서드는 Result 패턴을 사용하여 예외를 던지지 않습니다.
"""

import io
import json
import time
from datetime import datetime
from pathlib import Path
from threading import Lock
from typing import Callable, Dict, List, Optional

# RFS Framework 임포트
try:
    from rfs.core.result import Failure, Result, Success
    from rfs.hof.core import pipe, compose
    from rfs.hof.guard import guard
    from rfs.hof.collections import compact_map, first
except ImportError:
    # 폴백: 기본 Result 구현
    from dataclasses import dataclass
    from typing import Generic, TypeVar, Union

    T = TypeVar("T")
    E = TypeVar("E")

    @dataclass
    class Success(Generic[T]):
        """성공 결과"""

        value: T

        def is_success(self) -> bool:
            return True

        def unwrap(self) -> T:
            return self.value

        def map(self, func):
            """값 변환"""
            return Success(func(self.value))

        def bind(self, func):
            """Result 체이닝"""
            return func(self.value)

    @dataclass
    class Failure(Generic[E]):
        """실패 결과"""

        error: E

        def is_success(self) -> bool:
            return False

        def unwrap_error(self) -> E:
            return self.error

        def map(self, func):
            """실패는 그대로 전달"""
            return self

        def bind(self, func):
            """실패는 그대로 전달"""
            return self

    Result = Union[Success[T], Failure[E]]
    
    # HOF 폴백 구현
    def pipe(*funcs):
        """함수 합성 (왼쪽에서 오른쪽)"""
        def _pipe(value):
            result = value
            for func in funcs:
                result = func(result)
            return result
        return _pipe
    
    def compose(*funcs):
        """함수 합성 (오른쪽에서 왼쪽)"""
        return pipe(*reversed(funcs))
    
    def guard(condition, on_true, on_false):
        """조건에 따른 분기"""
        return on_true() if condition else on_false()
    
    def compact_map(func, items):
        """변환 후 None 제거"""
        return [result for item in items if (result := func(item)) is not None]
    
    def first(predicate, items):
        """조건을 만족하는 첫 번째 요소"""
        for item in items:
            if predicate(item):
                return item
        return None

import spacy
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from pypdf import PdfReader
from spacy.cli import download as spacy_download
from spacy.language import Language

from .entities import DocumentChunk, ProcessedDocument
from ..shared.logging import get_logger

# 로거 초기화
logger = get_logger(__name__)


# ============================================================================
# 파일 텍스트 추출 서비스 (단일책임 원칙 적용)
# ============================================================================

class FileTextExtractionService:
    """
    업로드된 파일에서 텍스트를 추출하는 서비스
    
    단일책임: 파일 형식 감지 및 적절한 추출기로 라우팅
    
    지원 형식:
        - 프레젠테이션: .ppt, .pptx
        - 문서: .doc, .docx
        - 엑셀: .xls, .xlsx
        - PDF: .pdf
        - 텍스트: .txt
        - 마크다운: .md, .markdown
        - JSON: .json
    """

    _EXTENSION_MAPPER: Dict[str, str] = {
        ".ppt": "presentation",
        ".pptx": "presentation",
        ".doc": "document",
        ".docx": "document",
        ".xls": "excel",
        ".xlsx": "excel",
        ".pdf": "pdf",
        ".txt": "text",
        ".md": "text",
        ".markdown": "text",
        ".json": "json",
    }

    def is_supported(self, filename: str) -> bool:
        """
        파일 확장자 지원 여부 확인
        
        Args:
            filename: 파일 이름
            
        Returns:
            bool: 지원 여부
        """
        suffix = Path(filename).suffix.lower()
        return suffix in self._EXTENSION_MAPPER

    def extract_text(self, file_bytes: bytes, filename: str) -> Result[str, str]:
        """
        파일에서 텍스트 추출 (HOF 패턴 적용)
        
        Args:
            file_bytes: 파일 바이너리
            filename: 원본 파일 이름

        Returns:
            Result[str, str]: 추출된 텍스트 또는 에러 메시지
        """
        start_time = time.time()
        
        logger.log_operation("text_extraction", "started", 
                           filename=filename, 
                           file_size=len(file_bytes))
        
        # 확장자 검증
        suffix = Path(filename).suffix.lower()
        handler_key = self._EXTENSION_MAPPER.get(suffix)
        
        if handler_key is None:
            logger.warning("지원하지 않는 파일 형식", 
                         filename=filename, 
                         extension=suffix)
            return Failure(f"지원하지 않는 파일 형식입니다: {suffix or 'unknown'}")
        
        # 핸들러 조회
        handler = getattr(self, f"_extract_{handler_key}", None)
        if handler is None:
            logger.error("파일 형식 처리기 없음", 
                       filename=filename, 
                       handler_key=handler_key)
            return Failure(f"해당 파일 형식 처리기가 구성되지 않았습니다: {suffix}")

        # 텍스트 추출 시도
        try:
            text_result = handler(file_bytes)
            clean_text = text_result.strip()
            
            # 빈 텍스트 검증
            if not clean_text:
                logger.warning("추출된 텍스트 없음", filename=filename)
                return Failure("파일에 추출 가능한 텍스트가 없습니다.")
            
            duration_ms = (time.time() - start_time) * 1000
            logger.log_performance("text_extraction", 
                                 duration_ms=duration_ms,
                                 filename=filename,
                                 text_length=len(clean_text),
                                 file_type=handler_key)
            
            logger.log_operation("text_extraction", "completed",
                               filename=filename,
                               text_length=len(clean_text))
            
            return Success(clean_text)
        except Exception as exc:
            duration_ms = (time.time() - start_time) * 1000
            logger.log_error_with_context(exc, "text_extraction",
                                        filename=filename,
                                        file_type=handler_key,
                                        duration_ms=duration_ms)
            return Failure(f"파일에서 텍스트를 추출하지 못했습니다: {str(exc)}")

    def _extract_document(self, file_bytes: bytes) -> str:
        """워드 문서(.doc, .docx)에서 텍스트 추출"""
        document = Document(io.BytesIO(file_bytes))

        # HOF: compact_map으로 문단별 텍스트 추출
        def extract_paragraph_text(paragraph) -> Optional[str]:
            text = paragraph.text.strip()
            return text if text else None

        paragraphs_text = compact_map(extract_paragraph_text, document.paragraphs)

        # 테이블 텍스트도 추출
        def extract_table_text(table) -> Optional[str]:
            table_lines = []
            for row in table.rows:
                row_cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if row_cells:
                    table_lines.append(" | ".join(row_cells))
            return "\n".join(table_lines) if table_lines else None

        tables_text = compact_map(extract_table_text, document.tables)

        # 문단과 테이블 텍스트 결합
        all_text = paragraphs_text + tables_text
        return "\n\n".join(all_text)

    def _extract_presentation(self, file_bytes: bytes) -> str:
        """파워포인트(.ppt, .pptx)에서 텍스트 추출"""
        presentation = Presentation(io.BytesIO(file_bytes))

        # HOF: compact_map으로 슬라이드별 텍스트 추출
        def extract_slide_text(slide_with_index) -> Optional[str]:
            slide_index, slide = slide_with_index
            slide_lines = [
                shape.text.strip()
                for shape in slide.shapes
                if hasattr(shape, "text") and shape.text.strip()
            ]
            if slide_lines:
                return f"[Slide {slide_index}]\n" + "\n".join(slide_lines)
            return None

        slides_text = compact_map(
            extract_slide_text,
            enumerate(presentation.slides, start=1)
        )
        return "\n\n".join(slides_text)

    def _extract_excel(self, file_bytes: bytes) -> str:
        """엑셀(.xls, .xlsx)에서 텍스트 추출"""
        workbook = load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
        
        # HOF: compact_map으로 시트별 텍스트 추출
        def extract_sheet_text(sheet) -> Optional[str]:
            rows = [
                " ".join(str(cell) for cell in row if cell is not None)
                for row in sheet.iter_rows(values_only=True)
            ]
            # 빈 행 제거
            non_empty_rows = [row for row in rows if row.strip()]
            if non_empty_rows:
                return f"[Sheet {sheet.title}]\n" + "\n".join(non_empty_rows)
            return None
        
        sheets_text = compact_map(extract_sheet_text, workbook.worksheets)
        return "\n\n".join(sheets_text)

    def _extract_pdf(self, file_bytes: bytes) -> str:
        """PDF에서 텍스트 추출"""
        reader = PdfReader(io.BytesIO(file_bytes))
        
        # HOF: compact_map으로 페이지별 텍스트 추출
        def extract_page_text(page_with_index) -> Optional[str]:
            page_index, page = page_with_index
            text = (page.extract_text() or "").strip()
            if text:
                return f"[Page {page_index}]\n{text}"
            return None
        
        pages = compact_map(
            extract_page_text,
            enumerate(reader.pages, start=1)
        )
        return "\n\n".join(pages)

    def _extract_text(self, file_bytes: bytes) -> str:
        """텍스트/마크다운 파일에서 텍스트 추출"""
        return self._decode_bytes(file_bytes)

    def _extract_json(self, file_bytes: bytes) -> str:
        """JSON 파일에서 텍스트 추출"""
        payload = json.loads(self._decode_bytes(file_bytes))
        return json.dumps(payload, ensure_ascii=False, indent=2)

    @staticmethod
    def _decode_bytes(file_bytes: bytes) -> str:
        """
        바이너리를 문자열로 디코딩
        
        Args:
            file_bytes: 바이너리 데이터
            
        Returns:
            str: 디코딩된 문자열
        """
        try:
            return file_bytes.decode("utf-8")
        except UnicodeDecodeError:
            return file_bytes.decode("latin-1", errors="ignore")


class DocumentChunkingService:
    """
    SpaCy 모델을 사용하여 문장 기준으로 텍스트를 청크로 분할하는 서비스
    
    단일책임: 텍스트 청킹만 담당 (모델 관리는 내부 구현)
    HOF 패턴: guard, compact_map 적용
    """

    _model_lock = Lock()
    _model: Language | None = None

    def __init__(self, model_name: str = "xx_sent_ud_sm", default_max_chars: int = 2000):
        """
        서비스 초기화
        
        Args:
            model_name: SpaCy 모델 이름
            default_max_chars: 기본 청크 최대 문자 수
        """
        self._model_name = model_name
        self._default_max_chars = default_max_chars

    def chunk_text(
        self, text: str, max_chars: Optional[int] = None
    ) -> Result[List[str], str]:
        """
        텍스트를 지정된 길이 기준으로 청크 분할 (HOF 패턴 적용)
        
        Args:
            text: 분할할 텍스트
            max_chars: 청크 최대 문자 수 (None이면 기본값 사용)
            
        Returns:
            Result[List[str], str]: 청크 리스트 또는 에러
        """
        start_time = time.time()
        
        logger.log_operation("text_chunking", "started",
                           text_length=len(text),
                           max_chars=max_chars or self._default_max_chars)
        
        # 입력 검증
        if not text or not text.strip():
            logger.warning("빈 텍스트 제공됨")
            return Failure("분할할 텍스트가 제공되지 않았습니다.")
        
        limit = max_chars or self._default_max_chars
        if limit <= 0:
            logger.error("잘못된 청크 크기", max_chars=limit)
            return Failure("청크 크기는 0보다 커야 합니다.")

        # SpaCy 모델로 문장 분리
        try:
            nlp = self._get_model()
            doc = nlp(text)
            
            # 문장 리스트를 청크로 변환
            sentences = [sent.text.strip() for sent in doc.sents if sent.text.strip()]
            chunks = self._build_chunks(sentences, limit)
            
            # 결과 검증
            if not chunks:
                logger.warning("청킹 결과 없음", text_length=len(text))
                return Failure("문장 분할 결과가 비어 있습니다.")
            
            duration_ms = (time.time() - start_time) * 1000
            logger.log_performance("text_chunking",
                                 duration_ms=duration_ms,
                                 text_length=len(text),
                                 sentence_count=len(sentences),
                                 chunk_count=len(chunks))
            
            logger.log_operation("text_chunking", "completed",
                               chunk_count=len(chunks),
                               total_chars=sum(len(c) for c in chunks))
            
            return Success(chunks)
        except Exception as exc:
            duration_ms = (time.time() - start_time) * 1000
            logger.log_error_with_context(exc, "text_chunking",
                                        text_length=len(text),
                                        duration_ms=duration_ms)
            return Failure(f"텍스트 청킹 실패: {str(exc)}")

    def _build_chunks(self, sentences: List[str], limit: int) -> List[str]:
        """
        문장 리스트를 청크로 조합 (함수형 접근)
        
        Args:
            sentences: 문장 리스트
            limit: 청크 최대 문자 수
            
        Returns:
            List[str]: 청크 리스트
        """
        chunks: List[str] = []
        current_chunk: List[str] = []
        current_length = 0

        for sentence in sentences:
            sentence_length = len(sentence)
            
            # 문장 자체가 limit보다 긴 경우: 강제 분할
            if sentence_length >= limit:
                # 현재 청크 마무리
                if current_chunk:
                    chunks.append(" ".join(current_chunk))
                    current_chunk = []
                    current_length = 0
                
                # 긴 문장을 limit 단위로 분할
                chunks.extend(self._split_long_sentence(sentence, limit))
                continue
            
            # 현재 청크에 문장 추가 가능 여부 확인
            projected_length = current_length + sentence_length + (1 if current_chunk else 0)
            
            if projected_length <= limit:
                # 현재 청크에 추가
                current_chunk.append(sentence)
                current_length = projected_length
            else:
                # 현재 청크 마무리 후 새 청크 시작
                if current_chunk:
                    chunks.append(" ".join(current_chunk))
                current_chunk = [sentence]
                current_length = sentence_length

        # 마지막 청크 추가
        if current_chunk:
            chunks.append(" ".join(current_chunk))

        return chunks

    @staticmethod
    def _split_long_sentence(sentence: str, limit: int) -> List[str]:
        """
        긴 문장을 limit 단위로 강제 분할
        
        Args:
            sentence: 분할할 문장
            limit: 분할 단위
            
        Returns:
            List[str]: 분할된 세그먼트 리스트
        """
        segments = []
        for start in range(0, len(sentence), limit):
            segment = sentence[start : start + limit].strip()
            if segment:
                segments.append(segment)
        return segments

    def _get_model(self) -> Language:
        """
        SpaCy 모델을 로드하거나 캐시된 모델을 반환
        
        Returns:
            Language: SpaCy 모델
        """
        with self._model_lock:
            if self._model is not None:
                return self._model

            try:
                self._model = spacy.load(self._model_name)
            except OSError:
                # 모델이 없으면 다운로드
                spacy_download(self._model_name)
                self._model = spacy.load(self._model_name)

            return self._model
